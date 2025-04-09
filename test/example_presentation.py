from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, Python-PPTX!"
subtitle.text = "This is a slide created with python-pptx."

bullet_slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(bullet_slide_layout)

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Bullet Points Slide"

tf = body_shape.text_frame
tf.text = "First bullet point"

p = tf.add_paragraph()
p.text = "Second bullet point"
p.level = 1

p = tf.add_paragraph()
p.text = "Third bullet point"
p.level = 2

prs.save("example_presentation.pptx")

