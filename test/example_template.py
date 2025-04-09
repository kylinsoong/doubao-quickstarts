from pptx import Presentation

# Load the existing template
prs = Presentation('template.pptx')

# Choose a slide layout from the template (e.g., title + content layout)
slide_layout = prs.slide_layouts[1]  # You can change this index based on your template

# Add a new slide using that layout
slide = prs.slides.add_slide(slide_layout)

# Set the title and content of the slide
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Quarterly Summary"
content.text = "• Revenue increased by 15%\n• New client acquisition\n• Expansion in APAC region"

# Save the new presentation
prs.save('custom_presentation.pptx')

